#!/usr/bin/env python3
"""pptx_inventory.py — read a (possibly hand-edited) .pptx into a structured
slide inventory JSON. Design A of the Talksmith reverse pipeline.

Given a rendered/edited deck and its Talk root, emit an ordered inventory of
every slide with:
  - classification: cover | agenda | section-divider | content
  - detected title, body blocks, speaker notes
  - embedded images, de-duplicated + filtered by a 5-tier ladder before staging
    to `reconcile/staging/` (see SKILL.md).

Only reading is performed — the inventory is consumed by reconstruct_md.py
to rebuild finalpptx.md.

Usage:
  python3 pptx_inventory.py <deck.pptx> --talk <talk_root> --style <strict|free-form>
                            [--stage-new] [--json | --human] [-o <inventory.json>]

Exit codes:
  0  inventory produced
  2  could not run (bad .pptx, missing talk root, unknown style)

Requires: python-pptx (`pip install python-pptx`).
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from dataclasses import dataclass, asdict, field
from pathlib import Path, PurePosixPath

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE

sys.path.insert(0, str(Path(__file__).resolve().parent.parent / "_shared"))
import _pptxlib as L

# --------------------------------------------------------------------------- #
# style palettes
# --------------------------------------------------------------------------- #

STYLE_PALETTE = {
    "strict": {
        "spine_fill": {"D8D4D4"},
        # Active-dot fills observed across real Cowork-rendered strict decks:
        #   DA1B2E, F33447 — base-template.pptx (hiperparametros-ai deck)
        #   E02020, F93939 — seguridad-governance-ai deck
        # All are variants of the brand red; the exact shade drifts with the
        # renderer version. Adding both here keeps the "exactly one active dot"
        # divider signal robust across template revisions.
        "dot_active_fill": {"DA1B2E", "F33447", "E02020", "F93939"},
        "dot_inactive_fill": {"F2EEEE"},
        "pill_fill": {"F9D2D6"},
        "callout_fill": {"F7BBC1", "B8E6F5"},
        "chrome": True,
    },
    "free-form": {
        "spine_fill": set(),
        "dot_active_fill": set(),
        "dot_inactive_fill": set(),
        "pill_fill": set(),
        "callout_fill": set(),
        "chrome": False,
    },
}

# Geometry thresholds (EMU). Widened to cover both the base-template.pptx
# (dots ≈ 250999²) and real Cowork-rendered agendas (dots ≈ 355790²).
SPINE_MAX_CX = 60_000
SPINE_MIN_CY = 3_000_000
DOT_MIN = 150_000
DOT_MAX = 500_000
DOT_SQUARE_TOL = 60_000
TOP_BAND_Y = 1_300_000

# Title-shape font-size threshold. Matches the original XML-side `sz >= 1700`
# (in hundredths of a point) → 17.0 pt.
TITLE_MIN_PT = 17.0


# --------------------------------------------------------------------------- #
# data model — unchanged from the pre-migration JSON schema
# --------------------------------------------------------------------------- #

@dataclass
class ImageRec:
    media: str
    sha256: str
    w: float | None
    h: float | None
    cls: str                    # content | template
    staged_path: str
    ordinal: int = 0


@dataclass
class SlideRec:
    order: int
    part: str
    classification: str
    classification_confidence: str
    signals: dict
    section_index: int | None
    section_name: str | None
    pill_text: str | None
    title: str
    body: list[dict]
    notes: str
    images: list[dict]
    comments: list[dict] = field(default_factory=list)


# --------------------------------------------------------------------------- #
# python-pptx shape helpers — small facades that also swallow the "shape has
# no fill / no auto_shape_type / no text_frame" AttributeError cases.
# --------------------------------------------------------------------------- #

def _shape_fill_hex(shape) -> str | None:
    """Solid-fill hex ('D8D4D4') or None."""
    try:
        fill = shape.fill
        if fill.type != 1:  # 1 = MSO_FILL.SOLID
            return None
        rgb = fill.fore_color.rgb
        return f"{rgb}".upper() if rgb else None
    except (AttributeError, ValueError, KeyError):
        return None


def _shape_is_roundrect(shape) -> bool:
    try:
        return shape.auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE
    except (AttributeError, ValueError, TypeError):
        return False


def _shape_geom(shape) -> tuple[tuple[int, int] | None, tuple[int, int] | None]:
    """Return ((left, top), (width, height)) in EMU, or (None, None)."""
    try:
        left, top = shape.left, shape.top
        width, height = shape.width, shape.height
    except (AttributeError, ValueError):
        return None, None
    off = (int(left), int(top)) if left is not None and top is not None else None
    ext = (int(width), int(height)) if width is not None and height is not None else None
    return off, ext


def _shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text or ""


def _max_font_pt(shape) -> float:
    """Largest run font size (in points) anywhere in the shape, or 0."""
    if not shape.has_text_frame:
        return 0.0
    best = 0.0
    for p in shape.text_frame.paragraphs:
        for r in p.runs:
            try:
                sz = r.font.size
            except AttributeError:
                sz = None
            if sz is None:
                continue
            try:
                pt = sz.pt
            except AttributeError:
                pt = float(sz) / 12700.0  # EMU → pt fallback
            if pt > best:
                best = pt
    return best


def _shape_paragraphs(shape) -> list[tuple[str, bool]]:
    """Return [(text, is_bulleted), ...] per paragraph in the shape."""
    if not shape.has_text_frame:
        return []
    out: list[tuple[str, bool]] = []
    for p in shape.text_frame.paragraphs:
        text = p.text.strip() if p.text else ""
        out.append((text, L.paragraph_is_bulleted(p)))
    return out


# --------------------------------------------------------------------------- #
# chrome / classification signals
# --------------------------------------------------------------------------- #

def _is_dot(shape) -> tuple[bool, str | None]:
    """(is_near_square_dot, fill_hex)."""
    if not _shape_is_roundrect(shape):
        return False, None
    _off, ext = _shape_geom(shape)
    if not ext:
        return False, None
    cx, cy = ext
    if not (DOT_MIN <= cx <= DOT_MAX and DOT_MIN <= cy <= DOT_MAX):
        return False, None
    if abs(cx - cy) > DOT_SQUARE_TOL:
        return False, None
    return True, _shape_fill_hex(shape)


def _is_spine(shape, spine_fills: set[str]) -> bool:
    if not _shape_is_roundrect(shape):
        return False
    _off, ext = _shape_geom(shape)
    if not ext:
        return False
    cx, cy = ext
    if not (cx < SPINE_MAX_CX and cy > SPINE_MIN_CY):
        return False
    fill = _shape_fill_hex(shape)
    return (fill in spine_fills) if spine_fills else True


def _is_pill(shape, pill_fills: set[str]) -> bool:
    if not pill_fills:
        return False
    if not _shape_is_roundrect(shape):
        return False
    return _shape_fill_hex(shape) in pill_fills


def _analyze_agenda_rows(slide, palette) -> list[dict]:
    """Walk shapes in document order; each dot opens a row, the next non-empty
    non-numeric text shape names it."""
    rows: list[dict] = []
    for shape in slide.shapes:
        is_dot, fill = _is_dot(shape)
        if is_dot:
            off, _ext = _shape_geom(shape)
            active = fill in palette["dot_active_fill"] if palette["dot_active_fill"] else False
            rows.append({"active": active, "title": None, "y": off[1] if off else 0})
            continue
        if rows and rows[-1]["title"] is None:
            text = _shape_text(shape).strip()
            if text and not text.isdigit() and not L.H1_AGENDA.match(f"# {text}"):
                rows[-1]["title"] = text
    return rows


def _detect_chrome(slide, palette) -> dict:
    agenda_title = False
    spine = False
    dot_count = 0
    active_dots = 0
    pill_text: str | None = None
    for shape in slide.shapes:
        text = _shape_text(shape).strip()
        if text and L.H1_AGENDA.match(f"# {text}"):
            agenda_title = True
        if _is_spine(shape, palette["spine_fill"]):
            spine = True
        is_dot, fill = _is_dot(shape)
        if is_dot:
            dot_count += 1
            if fill in palette["dot_active_fill"]:
                active_dots += 1
        if _is_pill(shape, palette["pill_fill"]) and pill_text is None:
            pill_text = text
    return {
        "agenda_title": agenda_title,
        "spine": spine,
        "dot_count": dot_count,
        "active_dots": active_dots,
        "pill_text": pill_text,
    }


# --------------------------------------------------------------------------- #
# title + body extraction
# --------------------------------------------------------------------------- #

def _extract_title(slide, palette) -> str:
    """Largest-font text shape in the top band, excluding the section pill."""
    best: tuple[float, int, str] | None = None
    for shape in slide.shapes:
        if _is_pill(shape, palette["pill_fill"]):
            continue
        text = _shape_text(shape).strip()
        if not text:
            continue
        off, _ext = _shape_geom(shape)
        y = off[1] if off else 0
        sz = _max_font_pt(shape)
        cand = (sz, -y, text)
        if y <= TOP_BAND_Y and (best is None or cand > best):
            best = cand
    if best:
        return best[2]
    # Fallback: largest anywhere
    best = None
    for shape in slide.shapes:
        if _is_pill(shape, palette["pill_fill"]):
            continue
        text = _shape_text(shape).strip()
        if not text:
            continue
        sz = _max_font_pt(shape)
        cand = (sz, 0, text)
        if best is None or cand > best:
            best = cand
    return best[2] if best else ""


def _extract_grid_tables(slide, palette, title_text: str) -> tuple[list[dict], set[int]]:
    """Detect tables rendered as arrangements of individual text shapes.

    Cowork's forward pipeline exports Markdown pipe-tables as a grid of
    text shapes at fixed x/y positions — NOT as native OOXML `<a:tbl>`.
    Without this reconstruction the table becomes a flat list of 40 text
    paragraphs on a slide, losing the row/column structure entirely.

    Detection rule:
      1. Collect body-worthy text shapes with their (x, y) origins.
      2. Cluster shapes into rows (same y within `Y_TOL` = 50k EMU).
      3. A "table" is a run of ≥3 consecutive rows all with the same
         column count AND the same column x-positions (within `X_TOL`).
         The ≥3-row threshold avoids false-positives on card-grid layouts
         (which are typically 1 heading-row + 1 body-row = 2 rows).

    Returns (tables, used_shape_ids). `_extract_body` skips shapes whose
    id() is in used_shape_ids so table cells aren't double-emitted.
    """
    Y_TOL = 50_000
    X_TOL = 100_000
    MIN_ROWS = 3

    candidates: list[tuple[int, int, str, int]] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or text == title_text:
            continue
        if _is_pill(shape, palette["pill_fill"]):
            continue
        if _is_spine(shape, palette["spine_fill"]):
            continue
        if _is_dot(shape)[0]:
            continue
        off, _ext = _shape_geom(shape)
        if not off:
            continue
        # `shape.shape_id` — OOXML `<p:cNvPr id="...">`, unique per slide.
        # NOTE: `id(shape)` is unstable across `slide.shapes` iterations because
        # python-pptx creates transient wrapper objects that get GC'd and
        # recycled — 8 shapes on a slide can share only 2 `id()` values,
        # collapsing the exclude-set and double-emitting table cells as body.
        candidates.append((off[1], off[0], text, shape.shape_id))

    if len(candidates) < MIN_ROWS * 2:  # need at least 3 rows × 2 cols
        return [], set()

    candidates.sort()
    rows: list[list[tuple[int, int, str, int]]] = []
    for y, x, text, sid in candidates:
        if rows and abs(y - rows[-1][0][0]) <= Y_TOL:
            rows[-1].append((y, x, text, sid))
        else:
            rows.append([(y, x, text, sid)])

    def _merge_cols(xs: list[int]) -> list[int]:
        """Collapse a sorted x-list into distinct column positions (within X_TOL)."""
        out: list[int] = []
        for x in sorted(xs):
            if out and abs(x - out[-1]) <= X_TOL:
                continue
            out.append(x)
        return out

    tables: list[dict] = []
    used: set[int] = set()
    i = 0
    while i < len(rows):
        if len(rows[i]) < 2:
            i += 1
            continue
        for r in rows[i:]:
            r.sort(key=lambda t: t[1])

        # Row-run extension. col_union is FIXED by the first row and never
        # grows — later rows may be a subset (missing cells for implicit
        # grouped columns like slide 13's ragged provider table) but must
        # never introduce columns not seen in row i. Growing the union would
        # let unrelated footnote rows attach as a phantom column (a
        # single-cell "Tendencia de fondo" at x=661475 fooled an earlier
        # implementation into treating it as column 0 of the table).
        col_union = _merge_cols([c[1] for c in rows[i]])
        j = i + 1
        while j < len(rows):
            row_xs = [c[1] for c in rows[j]]
            if not all(any(abs(x - cx) <= X_TOL for cx in col_union) for x in row_xs):
                break
            j += 1

        if j - i >= MIN_ROWS and len(col_union) >= 2:
            table_rows: list[list[str]] = []
            for r in rows[i:j]:
                cells = [""] * len(col_union)
                for _y, x, text, _sid in r:
                    for k, cx in enumerate(col_union):
                        if abs(x - cx) <= X_TOL:
                            cells[k] = text
                            break
                table_rows.append(cells)
            tables.append({"role": "table", "rows": table_rows})
            for r in rows[i:j]:
                for c in r:
                    used.add(c[3])
            i = j
        else:
            i += 1
    return tables, used


def _extract_body(slide, palette, title_text: str, exclude_ids: set[int] | None = None) -> list[dict]:
    """Body blocks (non-title, non-pill, non-chrome text shapes) in reading order.

    `exclude_ids` — `shape.shape_id`s consumed by `_extract_grid_tables` and
    emitted as table blocks; skipped here so their cell text isn't double-
    reported. Uses `shape.shape_id` (stable OOXML identifier) rather than
    `id(shape)` (unstable Python-object identity — see `_extract_grid_tables`).
    """
    exclude_ids = exclude_ids or set()
    blocks: list[tuple[int, int, dict]] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape.shape_id in exclude_ids:
            continue
        fill = _shape_fill_hex(shape)
        if _is_pill(shape, palette["pill_fill"]):
            continue
        if _is_spine(shape, palette["spine_fill"]):
            continue
        if _is_dot(shape)[0]:
            continue
        paras = _shape_paragraphs(shape)
        text = "\n".join(t for t, _b in paras if t).strip()
        if not text or text == title_text:
            continue
        off, _ext = _shape_geom(shape)
        y = off[1] if off else 0
        x = off[0] if off else 0
        if palette["callout_fill"] and fill in palette["callout_fill"]:
            role = "callout"
        elif sum(1 for t, b in paras if t and b) >= 2:
            role = "bullets"
        else:
            role = "paragraph"
        block = {"role": role, "fill": fill}
        if role == "bullets":
            block["items"] = [t for t, _b in paras if t]
        else:
            block["text"] = text
        blocks.append((y, x, block))
    blocks.sort(key=lambda t: (t[0], t[1]))
    return [b for _y, _x, b in blocks]


def _extract_tables(slide) -> list[dict]:
    out: list[dict] = []
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        rows: list[list[str]] = []
        for row in shape.table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        if rows:
            out.append({"role": "table", "rows": rows})
    return out


def _notes_for_slide(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    tf = slide.notes_slide.notes_text_frame
    return (tf.text or "").strip()


# --------------------------------------------------------------------------- #
# image de-duplication and classification (unchanged JSON output)
# --------------------------------------------------------------------------- #

def _index_images_dir(images_dir: Path) -> dict[str, str]:
    """Return {sha256: talk_rel_path} for every file under images/ (excluding
    staging, .ascii sidecars). Used by downstream diff for byte-identical
    fast-path."""
    by_hash: dict[str, str] = {}
    if not images_dir.is_dir():
        return by_hash
    talk_root = images_dir.parent
    for p in sorted(images_dir.rglob("*")):
        if not p.is_file() or p.name.endswith(".ascii"):
            continue
        h = L.sha256_file(p)
        if h and h not in by_hash:
            by_hash[h] = str(p.relative_to(talk_root))
    return by_hash


def _classify_and_stage(media_path: str, blob: bytes, images_dir: Path,
                        slide_order: int, ordinal: int) -> ImageRec:
    """Stage a content image to reconcile/staging/ and produce its ImageRec."""
    sha = L.sha256_bytes(blob)
    d = L.intrinsic_dims(media_path, blob)
    w, h = (d if d else (None, None))
    talk_root = images_dir.parent
    staging = talk_root / "reconcile" / "staging"
    staging.mkdir(parents=True, exist_ok=True)
    ext = PurePosixPath(media_path).suffix or ".png"
    dest = staging / f"slide{slide_order}-img{ordinal}{ext}"
    if not (dest.exists() and dest.read_bytes() == blob):
        dest.write_bytes(blob)
    return ImageRec(media=media_path, sha256=sha, w=w, h=h, cls="content",
                    staged_path=str(dest.relative_to(talk_root)),
                    ordinal=ordinal)


# --------------------------------------------------------------------------- #
# classification (cover / agenda / divider / content)
# --------------------------------------------------------------------------- #

def _classify(order: int, sig: dict, logo: bool, palette: dict,
              section_names: list) -> tuple[str, str, int | None]:
    if not palette["chrome"]:
        if order == 1:
            return "cover", "high", None
        return "content", "low", None
    is_agenda_chrome = (sig["agenda_title"] and sig["spine"] and sig["dot_count"] >= 2)
    if order == 1 and logo and not is_agenda_chrome:
        return "cover", "high", None
    if logo and not is_agenda_chrome and sig["dot_count"] == 0:
        return "cover", "medium", None
    if is_agenda_chrome:
        active = sig["active_dots"]
        sec_index = sig.get("_active_index")
        if active == 1:
            kind = "agenda" if order <= 2 else "section-divider"
            return kind, "high", sec_index
        return "section-divider", "low", sec_index
    return "content", "high", None


def _summarize(slides: list[SlideRec]) -> dict:
    def count(c): return sum(1 for s in slides if s.classification == c)
    return {
        "slides": len(slides),
        "cover": count("cover"),
        "agenda": count("agenda"),
        "dividers": count("section-divider"),
        "content": count("content"),
        "images_staged": sum(1 for s in slides
                             for im in s.images if im["cls"] == "content"),
        "comments": sum(len(s.comments) for s in slides),
    }


# --------------------------------------------------------------------------- #
# main build
# --------------------------------------------------------------------------- #

def build_inventory(pptx_path: str, talk_root: Path, style: str) -> dict:
    palette = STYLE_PALETTE[style]
    images_dir = talk_root / "images"
    by_hash = _index_images_dir(images_dir)

    try:
        prs = Presentation(pptx_path)
    except Exception as e:  # PackageNotFoundError, IsADirectoryError, etc.
        raise RuntimeError(f"cannot open {pptx_path}: {e}") from e

    warnings: list[str] = []

    # First pass: derive section_names from the first agenda-chrome slide.
    section_names: list[str] = []
    if palette["chrome"]:
        for slide in prs.slides:
            sig = _detect_chrome(slide, palette)
            if sig["agenda_title"] and sig["spine"] and sig["dot_count"] >= 2:
                rows = _analyze_agenda_rows(slide, palette)
                names = [r["title"] for r in rows if r["title"]]
                if names:
                    section_names = names
                    break

    # Second pass: classify, extract, stage.
    slides_out: list[SlideRec] = []
    current_section_index: int | None = None

    for order, slide in enumerate(prs.slides, start=1):
        # slide.part.partname is like '/ppt/slides/slide7.xml' — strip leading '/'.
        part_path = slide.part.partname.lstrip("/")

        sig = _detect_chrome(slide, palette)
        if palette["chrome"] and sig["agenda_title"] and sig["spine"] and sig["dot_count"] >= 2:
            rows = _analyze_agenda_rows(slide, palette)
            active_idx = next((i + 1 for i, r in enumerate(rows) if r["active"]), None)
            sig["_active_index"] = active_idx

        # Cover-logo detection: any picture on this slide whose raster media
        # path matches the known template basename regex (`image-1-*` etc.).
        # Package part-path — NOT `image.filename`, which is a collision-prone
        # synthetic basename.
        logo = False
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                fn = L.pic_raster_target(shape, slide.part) or ""
                if L.TEMPLATE_MEDIA_RE.search(fn) \
                        and PurePosixPath(fn).name.startswith("image-1-"):
                    logo = True
                    break

        classification, confidence, sec_index = _classify(
            order, sig, logo, palette, section_names)
        if classification in ("agenda", "section-divider"):
            current_section_index = sec_index
        elif classification == "content":
            sec_index = current_section_index

        # Reviewer comments are extracted for every slide type (cover, agenda,
        # divider, content). Even a cover-slide comment may carry review
        # feedback the presenter wants to see.
        comments = L.slide_comments(slide)

        if classification in ("cover", "agenda", "section-divider"):
            title = _extract_title(slide, palette) if classification == "cover" else ""
            body: list[dict] = []
            notes = ""
            images: list[dict] = []
        else:
            title = _extract_title(slide, palette)
            # Detect tables first (both native `<a:tbl>` AND grid-of-text-shapes
            # patterns), then subtract those shapes from the body pool so they
            # aren't double-emitted as loose paragraphs.
            grid_tables, table_shape_ids = _extract_grid_tables(slide, palette, title)
            body = _extract_body(slide, palette, title, exclude_ids=table_shape_ids)
            body += grid_tables
            body += _extract_tables(slide)
            notes = _notes_for_slide(slide)
            images = _extract_images(slide, images_dir, order)

        sec_name = None
        if sec_index and section_names and 1 <= sec_index <= len(section_names):
            sec_name = section_names[sec_index - 1]
        elif sig.get("pill_text"):
            sec_name = sig["pill_text"]

        slides_out.append(SlideRec(
            order=order, part=part_path, classification=classification,
            classification_confidence=confidence, signals=sig,
            section_index=sec_index, section_name=sec_name,
            pill_text=sig.get("pill_text"), title=title,
            body=body, notes=notes, images=images, comments=comments))

    return {
        "pptx": str(pptx_path),
        "talk_root": str(talk_root),
        "style": style,
        "section_names": section_names,
        "known_hashes": by_hash,
        "summary": _summarize(slides_out),
        "slides": [asdict(s) for s in slides_out],
        "warnings": warnings,
    }


def _extract_images(slide, images_dir: Path, order: int) -> list[dict]:
    """Run the 5-tier filter ladder and stage the survivors.

    Tier order matches the pre-migration behavior exactly:
      1. TEMPLATE_MEDIA_RE (icon-*, image-1-*)
      2. same media path reused ≥2× on this slide → chrome
      3. SVG icon (colorable-icon / data-icon / tiny render / zero <text>)
      4. PNG-only icon by size (max ≤128)
      5. raster floor: len(blob) < 10 KB
    """
    part = slide.part
    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]

    # Pre-scan: count how many times each raster path appears on this slide.
    # Package part-path (via rels) — NOT `image.filename`, which is a
    # synthetic basename that collides across slides.
    media_use_count: dict[str, int] = {}
    for pic in pictures:
        fn = L.pic_raster_target(pic, part) or ""
        if fn.lower().endswith((".png", ".jpg", ".jpeg")):
            media_use_count[fn] = media_use_count.get(fn, 0) + 1

    images: list[dict] = []
    content_ordinal = 0
    for pic in pictures:
        raster_path = L.pic_raster_target(pic, part)
        # SVG twin (via lxml drop-through) — icon-content sniffing + the SVG-only fallback.
        svg_path = L.pic_svg_target(pic, part)
        svg_blob = L.pic_svg_blob(pic, part) if svg_path else None

        # Prefer the raster; fall back to SVG when the pic carries only an svgBlip
        # (an embedded raster is absent — e.g. an SVG-only or link-only picture).
        picked = raster_path if raster_path else svg_path
        if not picked:
            continue
        img = None
        if raster_path:
            try:
                img = pic.image
            except (AttributeError, KeyError):
                # Raster declared but unreadable (external r:link, missing part) —
                # fall back to the SVG twin if one exists, else drop the pic.
                if not svg_blob:
                    continue
                picked = svg_path

        # Filter 1: template basename regex.
        if L.TEMPLATE_MEDIA_RE.search(picked):
            continue

        # Filter 2: same raster reused within the slide → card-layout chrome.
        if raster_path and media_use_count.get(raster_path, 0) > 1:
            continue

        # Filter 3: SVG content sniff (runs on the SVG twin bytes if present).
        if svg_blob and L.is_svg_icon(svg_blob):
            continue

        # Load the picked blob.
        if picked == raster_path and img is not None:
            blob = img.blob
        else:
            blob = svg_blob or (img.blob if img is not None else None)
        if blob is None:
            continue

        # Filter 4: PNG-only icon by intrinsic size.
        if picked.lower().endswith(".png") and L.is_png_icon(blob):
            continue

        # Filter 5: sub-10-KB raster floor.
        if picked.lower().endswith((".png", ".jpg", ".jpeg")) and len(blob) < 10240:
            continue

        content_ordinal += 1
        rec = _classify_and_stage(picked, blob, images_dir, order, content_ordinal)
        images.append(asdict(rec))
    return images


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #

def _human(inv: dict) -> str:
    s = inv["summary"]
    lines = [
        f"inventory: {s['slides']} slides "
        f"({s['cover']} cover, {s['agenda']} agenda, {s['dividers']} dividers, "
        f"{s['content']} content) — {s['images_staged']} content image(s) "
        f"staged under reconcile/staging/"
    ]
    if inv["section_names"]:
        lines.append(f"  sections: {', '.join(inv['section_names'])}")
    for sl in inv["slides"]:
        tag = sl["classification"][:4]
        extra = f' "{sl["title"]}"' if sl["title"] else ""
        conf = "" if sl["classification_confidence"] == "high" \
            else f" ⚠{sl['classification_confidence']}"
        imgs = f" [{len(sl['images'])} img]" if sl["images"] else ""
        lines.append(f"  {sl['order']:>2} {tag:<4}{conf} sec={sl['section_index']}{extra}{imgs}")
    for w in inv["warnings"]:
        lines.append(f"  ⚠ {w}")
    return "\n".join(lines)


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("pptx")
    ap.add_argument("--talk", required=True, help="Talk root, e.g. talks/<Talk>")
    ap.add_argument("--style", required=True, choices=sorted(STYLE_PALETTE))
    ap.add_argument("--stage-new", action="store_true",
                    help="(deprecated no-op — every content image is always staged "
                         "under reconcile/staging/)")
    ap.add_argument("--json", action="store_true")
    ap.add_argument("--human", action="store_true")
    ap.add_argument("-o", "--output",
                    help="write inventory JSON here "
                         "(default: <talk>/reconcile/finalpptx.inventory.json)")
    args = ap.parse_args(argv)

    talk_root = Path(args.talk)
    if not talk_root.is_dir():
        print(f"failed: talk root not found: {talk_root}", file=sys.stderr)
        return 2
    if not Path(args.pptx).is_file():
        print(f"failed: pptx not found: {args.pptx}", file=sys.stderr)
        return 2

    try:
        inv = build_inventory(args.pptx, talk_root, args.style)
    except RuntimeError as e:
        print(f"failed: {e}", file=sys.stderr)
        return 2

    out_path = Path(args.output) if args.output else talk_root / "reconcile" / "finalpptx.inventory.json"
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text(json.dumps(inv, indent=2, ensure_ascii=False) + "\n",
                        encoding="utf-8")

    if args.human:
        print(_human(inv))
    else:
        print(json.dumps(inv, indent=2, ensure_ascii=False))
    print(f"\ninventory: wrote {out_path}", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
