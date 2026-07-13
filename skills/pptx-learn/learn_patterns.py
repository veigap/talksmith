"""Learn strict styling / distribution / positioning patterns from human edits.

Strict-only. Talksmith renders a deck; the presenter opens it in Keynote/PowerPoint
and hand-corrects positions, sizes, fonts, fills. Those corrections are signal: they
say what the generator *should* have done. This script extracts a per-shape geometry
inventory from a deck, diffs the human-edited deck against the as-generated baseline,
and mines the **recurring** deltas into candidate conformance patterns — which a human
then curates into `config/pptx-styles/pptx-strict/conformance-patterns.md`.

Two subcommands:

    learn_patterns.py inventory <deck.pptx> [-o geom.json]
        Extract a geometry inventory. Run at render time on the as-generated deck to
        snapshot the baseline (survives in-place editing), and on the edited deck.

    learn_patterns.py diff --baseline <base.pptx|geom.json> --edited <edit.pptx|geom.json> \
        [-o candidates.json] [--min-recur 3] [--min-move-emu 45720]
        Match slides (by title, then order) and shapes (by role, then nearest position),
        compute per-shape deltas, and aggregate the recurring ones into candidate
        patterns (move / resize / refont / refill). Emits candidates in the shape
        `conformance-patterns.md` uses, so promotion is copy-paste.

Requires python-pptx. stdlib otherwise.
"""

from __future__ import annotations

import argparse
import json
import re
import statistics
import sys
from pathlib import Path

EMU_PER_INCH = 914400

# A shape moved/resized by less than this (EMU) is noise, not a pattern. 45720 = 0.05in.
DEFAULT_MIN_MOVE_EMU = 45720
DEFAULT_MIN_RECUR = 3

_PILL_FILLS = {"F9D2D6"}
_CALLOUT_FILLS = {"F7BBC1", "B8E6F5"}


def _norm_title(s: str) -> str:
    return re.sub(r"\s+", " ", (s or "").strip().lower())


def _fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type is not None and fill.fore_color and fill.fore_color.type is not None:
            rgb = fill.fore_color.rgb
            return str(rgb) if rgb is not None else None
    except (AttributeError, TypeError, ValueError):
        return None
    return None


def _max_font_pt(shape) -> float:
    best = 0.0
    if not getattr(shape, "has_text_frame", False):
        return best
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            sz = run.font.size
            if sz is not None:
                best = max(best, sz.pt)
    return best


def _text_snippet(shape, n: int = 40) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    t = re.sub(r"\s+", " ", shape.text_frame.text.strip())
    return t[:n]


def _geom(shape):
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
    except (AttributeError, ValueError):
        return None, None
    off = [int(l), int(t)] if l is not None and t is not None else None
    ext = [int(w), int(h)] if w is not None and h is not None else None
    return off, ext


def _kind(shape) -> str:
    st = str(getattr(shape, "shape_type", "")).lower()
    if "picture" in st:
        return "pic"
    if "table" in st or getattr(shape, "has_table", False):
        return "table"
    if "placeholder" in st:
        return "placeholder"
    if getattr(shape, "has_text_frame", False):
        return "textbox"
    return "autoshape"


def _role(shape, kind, off, ext, font_pt, fill, title_id) -> str:
    """Coarse role guess from position/size/fill — enough to match shapes across decks."""
    if shape.shape_id == title_id:
        return "title"
    if kind == "pic":
        # Small picture in a corner → icon; otherwise a content image.
        if ext and max(ext) <= EMU_PER_INCH // 2:
            return "icon"
        return "image"
    if fill and fill.upper() in _PILL_FILLS:
        return "pill"
    if fill and fill.upper() in _CALLOUT_FILLS:
        return "callout"
    if kind in ("textbox", "placeholder") and _text_snippet(shape):
        return "body"
    return "other"


def _classify_slide(order, shapes) -> str:
    if order == 1:
        return "cover"
    texts = [s for s in shapes if s["text"]]
    if any("agenda" in s["text"].lower() for s in shapes[:3]):
        return "agenda"
    # Divider: one dominant title, essentially no body text shapes.
    body = [s for s in shapes if s["role"] == "body"]
    if len(body) == 0 and any(s["role"] == "title" for s in shapes):
        return "section-divider"
    return "content-slide"


def _extract_title_id(slide) -> int | None:
    """Largest-font text shape in the top third — the slide title."""
    best_id, best_sz = None, 0.0
    for shape in slide.shapes:
        off, _ = _geom(shape)
        if not off or off[1] > 1_800_000:  # top ~third of a 5.6M-EMU-tall slide
            continue
        sz = _max_font_pt(shape)
        if sz > best_sz:
            best_sz, best_id = sz, shape.shape_id
    return best_id


def inventory_deck(pptx_path: str) -> dict:
    from pptx import Presentation

    prs = Presentation(pptx_path)
    slides_out = []
    for order, slide in enumerate(prs.slides, start=1):
        title_id = _extract_title_id(slide)
        shapes = []
        title_text = ""
        for shape in slide.shapes:
            off, ext = _geom(shape)
            kind = _kind(shape)
            fill = _fill_hex(shape)
            font_pt = _max_font_pt(shape)
            role = _role(shape, kind, off, ext, font_pt, fill, title_id)
            text = _text_snippet(shape)
            if shape.shape_id == title_id:
                title_text = text
            shapes.append({
                "id": int(shape.shape_id),
                "kind": kind,
                "role": role,
                "off": off,
                "ext": ext,
                "font_pt": round(font_pt, 1),
                "fill": fill,
                "text": text,
            })
        klass = _classify_slide(order, shapes)
        slides_out.append({
            "order": order,
            "title": title_text,
            "class": klass,
            "shapes": shapes,
        })
    return {"pptx": str(pptx_path), "slides": slides_out}


def _load_geom(path: str) -> dict:
    p = Path(path)
    if p.suffix.lower() == ".json":
        return json.loads(p.read_text(encoding="utf-8"))
    return inventory_deck(path)


def _match_slides(base: list[dict], edited: list[dict]) -> list[tuple[dict, dict]]:
    """Pair slides by normalized title; unmatched fall back to order alignment."""
    pairs = []
    by_title = {}
    for s in base:
        by_title.setdefault(_norm_title(s["title"]), []).append(s)
    used = set()
    for e in edited:
        key = _norm_title(e["title"])
        cands = by_title.get(key)
        if key and cands:
            b = cands.pop(0)
            pairs.append((b, e))
            used.add(b["order"])
    # Fallback: align leftover slides by order.
    leftover_base = [s for s in base if s["order"] not in used]
    matched_edited_orders = {e["order"] for _, e in pairs}
    leftover_edited = [e for e in edited if e["order"] not in matched_edited_orders]
    for b, e in zip(leftover_base, leftover_edited):
        pairs.append((b, e))
    return pairs


def _match_shapes(b_shapes: list[dict], e_shapes: list[dict]) -> list[tuple[dict, dict]]:
    """Pair shapes: same role, then nearest baseline position; unique 1:1."""
    pairs = []
    used_e = set()
    for b in b_shapes:
        best, best_d = None, None
        for i, e in enumerate(e_shapes):
            if i in used_e or e["role"] != b["role"]:
                continue
            if not b["off"] or not e["off"]:
                d = 0 if (not b["off"] and not e["off"]) else 1 << 60
            else:
                d = abs(b["off"][0] - e["off"][0]) + abs(b["off"][1] - e["off"][1])
            if best_d is None or d < best_d:
                best_d, best = d, i
        if best is not None:
            used_e.add(best)
            pairs.append((b, e_shapes[best]))
    return pairs


def _deltas(pairs_slides: list[tuple[dict, dict]], min_move: int) -> list[dict]:
    """Per matched shape, emit a delta record if it moved/resized/refonted/refilled."""
    out = []
    for b_slide, e_slide in pairs_slides:
        for b, e in _match_shapes(b_slide["shapes"], e_slide["shapes"]):
            rec = {"class": e_slide["class"], "role": b["role"],
                   "slide": e_slide["order"], "title": e_slide["title"]}
            if b["off"] and e["off"]:
                dx, dy = e["off"][0] - b["off"][0], e["off"][1] - b["off"][1]
                if abs(dx) >= min_move or abs(dy) >= min_move:
                    out.append({**rec, "type": "move", "dx": dx, "dy": dy})
            if b["ext"] and e["ext"]:
                dw, dh = e["ext"][0] - b["ext"][0], e["ext"][1] - b["ext"][1]
                if abs(dw) >= min_move or abs(dh) >= min_move:
                    out.append({**rec, "type": "resize", "dw": dw, "dh": dh})
            if b["font_pt"] and e["font_pt"] and abs(e["font_pt"] - b["font_pt"]) >= 0.5:
                out.append({**rec, "type": "refont", "d_pt": round(e["font_pt"] - b["font_pt"], 1)})
            if b["fill"] and e["fill"] and b["fill"] != e["fill"]:
                out.append({**rec, "type": "refill", "from": b["fill"], "to": e["fill"]})
    return out


def _median_int(xs):
    return int(round(statistics.median(xs)))


def aggregate(deltas: list[dict], min_recur: int) -> list[dict]:
    """Group deltas by (class, role, type[, direction]) and keep the recurring ones."""
    groups: dict[tuple, list[dict]] = {}
    for d in deltas:
        if d["type"] == "move":
            key = (d["class"], d["role"], "move",
                   (d["dx"] > 0) - (d["dx"] < 0), (d["dy"] > 0) - (d["dy"] < 0))
        elif d["type"] == "resize":
            key = (d["class"], d["role"], "resize",
                   (d["dw"] > 0) - (d["dw"] < 0), (d["dh"] > 0) - (d["dh"] < 0))
        elif d["type"] == "refont":
            key = (d["class"], d["role"], "refont", (d["d_pt"] > 0) - (d["d_pt"] < 0))
        else:  # refill
            key = (d["class"], d["role"], "refill", d["from"], d["to"])
        groups.setdefault(key, []).append(d)

    candidates = []
    for key, ds in groups.items():
        if len(ds) < min_recur:
            continue
        klass, role, typ = key[0], key[1], key[2]
        slides = sorted(x["slide"] for x in ds)
        cand = {"applies_to": klass, "role": role, "type": typ,
                "count": len(ds), "slides": slides}
        if typ == "move":
            cand["dx_emu"] = _median_int([x["dx"] for x in ds])
            cand["dy_emu"] = _median_int([x["dy"] for x in ds])
            cand["summary"] = (f"{role} on {klass} moved "
                               f"{cand['dx_emu']/EMU_PER_INCH:+.2f}in x, "
                               f"{cand['dy_emu']/EMU_PER_INCH:+.2f}in y "
                               f"({len(ds)} slides)")
        elif typ == "resize":
            cand["dw_emu"] = _median_int([x["dw"] for x in ds])
            cand["dh_emu"] = _median_int([x["dh"] for x in ds])
            cand["summary"] = (f"{role} on {klass} resized "
                               f"{cand['dw_emu']/EMU_PER_INCH:+.2f}in w, "
                               f"{cand['dh_emu']/EMU_PER_INCH:+.2f}in h ({len(ds)} slides)")
        elif typ == "refont":
            cand["d_pt"] = round(statistics.median([x["d_pt"] for x in ds]), 1)
            cand["summary"] = f"{role} on {klass} font {cand['d_pt']:+}pt ({len(ds)} slides)"
        else:
            cand["from"], cand["to"] = key[3], key[4]
            cand["summary"] = f"{role} on {klass} fill {key[3]}→{key[4]} ({len(ds)} slides)"
        candidates.append(cand)
    candidates.sort(key=lambda c: c["count"], reverse=True)
    return candidates


def _cmd_inventory(args) -> int:
    inv = inventory_deck(args.deck)
    out = json.dumps(inv, indent=2, ensure_ascii=False)
    if args.output:
        Path(args.output).parent.mkdir(parents=True, exist_ok=True)
        Path(args.output).write_text(out + "\n", encoding="utf-8")
        print(f"[pptx-learn] geometry inventory: {len(inv['slides'])} slides → {args.output}",
              file=sys.stderr)
    else:
        sys.stdout.write(out + "\n")
    return 0


def _cmd_diff(args) -> int:
    base = _load_geom(args.baseline)
    edited = _load_geom(args.edited)
    pairs = _match_slides(base["slides"], edited["slides"])
    deltas = _deltas(pairs, args.min_move_emu)
    candidates = aggregate(deltas, args.min_recur)
    doc = {
        "baseline": base.get("pptx"),
        "edited": edited.get("pptx"),
        "slides_matched": len(pairs),
        "raw_deltas": len(deltas),
        "min_recur": args.min_recur,
        "candidates": candidates,
    }
    out = json.dumps(doc, indent=2, ensure_ascii=False)
    if args.output:
        Path(args.output).parent.mkdir(parents=True, exist_ok=True)
        Path(args.output).write_text(out + "\n", encoding="utf-8")
    else:
        sys.stdout.write(out + "\n")
    print(f"[pptx-learn] {len(pairs)} slides matched, {len(deltas)} deltas → "
          f"{len(candidates)} candidate pattern(s)", file=sys.stderr)
    for c in candidates:
        print(f"  • {c['summary']}", file=sys.stderr)
    return 0


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    sub = ap.add_subparsers(dest="cmd", required=True)

    p_inv = sub.add_parser("inventory", help="extract a geometry inventory from a deck")
    p_inv.add_argument("deck")
    p_inv.add_argument("-o", "--output")
    p_inv.set_defaults(func=_cmd_inventory)

    p_diff = sub.add_parser("diff", help="diff edited vs baseline → candidate patterns")
    p_diff.add_argument("--baseline", required=True, help="as-generated deck.pptx or geom.json")
    p_diff.add_argument("--edited", required=True, help="human-edited deck.pptx or geom.json")
    p_diff.add_argument("-o", "--output")
    p_diff.add_argument("--min-recur", type=int, default=DEFAULT_MIN_RECUR)
    p_diff.add_argument("--min-move-emu", type=int, default=DEFAULT_MIN_MOVE_EMU)
    p_diff.set_defaults(func=_cmd_diff)

    args = ap.parse_args(argv)
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("failed: python-pptx is required (`pip install python-pptx`)", file=sys.stderr)
        return 3
    return args.func(args)


if __name__ == "__main__":
    raise SystemExit(main())
