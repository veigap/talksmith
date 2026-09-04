#!/usr/bin/env python3
"""Derive an editable `.pptx` from the rendered HTML deck.

The HTML deck is the single source of truth for how a slide looks. This does not re-implement
that look: it *measures* it. Chrome lays the deck out, `templates/html/harvest.js` walks the
laid-out DOM and reports a display list, and this rebuilds that display list as native
PowerPoint shapes — real text boxes with real, editable runs, not pictures of slides.

Why that is worth the machinery: the previous `.pptx` path was a 1143-line prose spec restating
every geometry rule in EMU for an LLM to follow, plus five OOXML audits to check it had. A
change to a card's corner radius had to be made twice, in two languages, and verified by a third.
Here there is one layout engine, and the exporter cannot disagree with it.

The coordinate mapping is exact, which is why this is faithful rather than approximate: the deck
is laid out at 1280x720 CSS px and a 16:9 PowerPoint slide is 12192000x6858000 EMU, so one CSS
pixel is exactly 9525 EMU on both axes and 0.75 pt of type. No rounding drift accumulates.

Fonts: a `.pptx` cannot portably embed one, so the export names families every machine has. The
substitution happens in the BROWSER, before layout (`_EXPORT_FONT_CSS` in html_style.py) — so the
fit pass binary-searches against Arial's metrics and the geometry harvested *is* Arial geometry.
Measuring in one font and naming another is the bug that would otherwise make every wrapped
block break in the wrong place; here it cannot arise, because this writes back the family the
harvest reported.

Usage:
    python3 export_pptx.py --talk talks/<Talk> [--style <skin>] [--theme light|dark]
    python3 export_pptx.py --deck path/to/index.html -o out.pptx
"""

from __future__ import annotations

import argparse
import base64
import io
import json
import re
import sys
from pathlib import Path

_HERE = Path(__file__).resolve().parent
sys.path.insert(0, str(_HERE))

import _chrome  # noqa: E402

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Emu, Pt
except ImportError:  # pragma: no cover
    sys.exit("[pptx] python-pptx is required: python3 -m pip install python-pptx")

try:
    from PIL import Image
except ImportError:  # pragma: no cover
    Image = None

# 12192000 / 1280 == 6858000 / 720 == 9525, exactly. The whole export rests on this being an
# integer on both axes: no scaling constant is ever guessed and no rounding drift accumulates.
EMU_PX = 9525
PT_PX = 0.75
SLIDE_W, SLIDE_H = 12192000, 6858000

# SVG rasterization. Icons paint at ~74px, so 3x gives a crisp 222px glyph; the cap keeps a
# full-bleed diagram from becoming a 3840px RGBA slab (four of those were 90% of an early
# export's weight). 2048px across a 13.3in slide is ~154 DPI — past what a projector or a laser
# printer resolves.
SVG_SCALE = 3
SVG_MAX = 2048
SHEET_W, SHEET_H = 2048, 4096


def E(v: float) -> Emu:
    return Emu(int(round(v * EMU_PX)))


def P(v: float) -> Pt:
    return Pt(round(v * PT_PX, 2))


def rgb(c) -> RGBColor | None:
    if not c:
        return None
    return RGBColor(int(c["r"]), int(c["g"]), int(c["b"]))


def alpha_of(c, extra: float = 1.0) -> float:
    return max(0.0, min(1.0, float(c.get("a", 1)) * extra)) if c else 0.0


def _set_alpha(fill_elm, a: float) -> None:
    """python-pptx has no translucency API. A fill is opaque unless its colour carries an
    `<a:alpha>`, and the deck uses translucency in real places (struck-through text at .55, the
    quiz explanation at .78, list bullets at .45), so drop the element in by hand."""
    if a >= 0.999:
        return
    clr = fill_elm.find(qn("a:srgbClr"))
    if clr is None:
        return
    el = clr.makeelement(qn("a:alpha"), {"val": str(int(round(a * 100000)))})
    clr.append(el)


# ---------------------------------------------------------------------------------------------
# Harvest
# ---------------------------------------------------------------------------------------------

def harvest(deck: Path, style: str = "", theme: str = "", timeout: int = 180) -> dict:
    """Drive the browser once and read the display list back out of the dumped DOM.

    The harvest rides `?print-pdf` because that is the only view where every slide is laid out at
    full size simultaneously; `export=pptx` both arms the harvester and stands the auto-print
    down. Theme and skin travel as URL params exactly the way the deck's own PDF button passes
    them, so the export is a snapshot of a chosen look rather than of whatever was last toggled.
    """
    q = "print-pdf&export=pptx"
    if theme:
        q += "&deck-theme=" + theme
    if style and style != "default":
        q += "&deck-style=" + style
    url = _chrome.file_url(deck, q)
    dom = _chrome.run(["--virtual-time-budget=%d" % (timeout * 1000), "--dump-dom"],
                      url, timeout=timeout).decode("utf-8", "replace")

    err = re.search(r'<script type="application/json" id="deck-harvest-error">(.*?)</script>',
                    dom, re.S)
    if err:
        d = json.loads(err.group(1).replace("<\\/", "</"))
        raise RuntimeError("the harvester failed inside the browser: %s\n%s"
                           % (d.get("error"), d.get("stack", "")[:600]))
    m = re.search(r'<script type="application/json" id="deck-harvest">(.*?)</script>', dom, re.S)
    ready = re.search(r'data-harvest-ready="(\d+)"', dom)
    if not m or not ready:
        raise RuntimeError(
            "the harvest did not complete within %ds.\n"
            "  The deck must be built by a current build_html.py — one rendered before the export\n"
            "  path existed carries no harvester. Re-render it and try again." % timeout)
    data = json.loads(m.group(1).replace("<\\/", "</"))
    if len(data["slides"]) != int(ready.group(1)):
        raise RuntimeError("harvest reported %s slides but carries %d"
                           % (ready.group(1), len(data["slides"])))
    return data


# ---------------------------------------------------------------------------------------------
# Pictures
# ---------------------------------------------------------------------------------------------

def _sprite_pages(items: list[tuple[str, int, int]]) -> list[list[tuple]]:
    """Shelf-pack the SVGs into sheets. Rasterizing one image per browser launch costs ~2s each
    and a deck carries ~120 of them; packing them into a couple of sheets turns two minutes into
    two seconds. The layout is computed here, not read back from the page, so the slicing
    coordinates are known without a second round trip."""
    pages, page, x, y, shelf = [], [], 0, 0, 0
    for key, w, h in items:
        w, h = max(1, min(w, SHEET_W)), max(1, h)
        if x + w > SHEET_W:
            x, y, shelf = 0, y + shelf, 0
        if y + h > SHEET_H and page:
            pages.append(page)
            page, x, y, shelf = [], 0, 0, 0
        page.append((key, x, y, w, h))
        x += w
        shelf = max(shelf, h)
    if page:
        pages.append(page)
    return pages


def rasterize_svgs(svgs: dict[str, str], sizes: dict[str, tuple[float, float]],
                   workdir: Path) -> dict[str, bytes]:
    """Rasterize every inline SVG through the same browser that drew it.

    Inline SVG is how the deck carries both its icon set and every ASCII diagram, and a `.pptx`
    has no vector import path worth relying on. Chrome is already the renderer of record, so
    using it here means an icon looks in PowerPoint exactly as it looked on screen — no second
    SVG implementation, and no `cairosvg` dependency that may not be installed.
    """
    if not svgs:
        return {}
    if Image is None:
        sys.stderr.write("[pptx] Pillow missing — diagrams and icons will be skipped\n")
        return {}

    items = []
    for k in svgs:
        w, h = sizes.get(k, (64, 64))
        f = min(SVG_SCALE, SVG_MAX / max(w, h, 1))
        items.append((k, max(1, int(round(w * f))), max(1, int(round(h * f)))))
    out: dict[str, bytes] = {}

    for n, page in enumerate(_sprite_pages(items)):
        pw = max(x + w for _, x, _, w, _ in page)
        ph = max(y + h for _, _, y, _, h in page)
        cells = "".join(
            '<div style="position:absolute;left:%dpx;top:%dpx;width:%dpx;height:%dpx">%s</div>'
            % (x, y, w, h, svgs[k]) for k, x, y, w, h in page)
        # `svg{width:100%;height:100%}` is the point of the wrapper: `_svg()` strips width/height
        # and keeps the viBox, so each glyph scales to whatever cell it is given.
        html = ('<!doctype html><meta charset="utf-8">'
                '<style>html,body{margin:0;background:transparent}'
                'div>svg{width:100%;height:100%;display:block}</style>' + cells)
        src = workdir / ("sprite-%d.html" % n)
        png = workdir / ("sprite-%d.png" % n)
        src.write_text(html, encoding="utf-8")
        _chrome.run(["--screenshot=%s" % png, "--window-size=%d,%d" % (pw, ph),
                     "--force-device-scale-factor=1", "--default-background-color=00000000",
                     "--virtual-time-budget=10000"], _chrome.file_url(src), timeout=120)
        if not png.exists():
            sys.stderr.write("[pptx] sprite sheet %d did not render\n" % n)
            continue
        sheet = Image.open(png).convert("RGBA")
        for k, x, y, w, h in page:
            tile = sheet.crop((x, y, min(x + w, sheet.width), min(y + h, sheet.height)))
            out[k] = _encode(tile)
    return out


def _encode(im) -> bytes:
    """Encode one rasterized graphic, choosing the format from what the graphic actually is.

    Both branches earn their place. Icons and ASCII diagrams are flat line art whose only extra
    shades are anti-aliasing, so reducing them to 256 colours makes the PNG compress ~20x better
    at no visible cost. But an atmospheric aside is a smooth gradient, and quantizing one bands
    it visibly, so anything with real tonal range keeps its full colour and, when it carries no
    transparency, goes to JPEG — which is what gradients compress well as.

    **Reduced-colour art is still written as RGBA, never as a palette PNG.** A palette image
    stores transparency as a `tRNS` table rather than a real alpha channel, which is the kind of
    PNG variant that trips importers; and for the icons that dominate this deck the RGBA file is
    actually the *smaller* of the two anyway. Likewise the JPEGs are baseline, not progressive:
    the saving is negligible and the compatibility is not.
    """
    buf = io.BytesIO()
    try:
        rgba = im.convert("RGBA")
        if rgba.getcolors(4096) is not None:                     # flat line art
            rgba.quantize(colors=256, method=Image.Quantize.FASTOCTREE).convert("RGBA").save(
                buf, "PNG", optimize=True)
            return buf.getvalue()
        if rgba.getchannel("A").getextrema()[0] >= 250:          # opaque tonal image
            rgba.convert("RGB").save(buf, "JPEG", quality=82, optimize=True, progressive=False)
            return buf.getvalue()
    except Exception:
        pass
    buf = io.BytesIO()
    im.save(buf, "PNG", optimize=True)
    return buf.getvalue()


_DATA_URI = re.compile(r"^data:([^;,]+)(;base64)?,(.*)$", re.S)


def decode_img(node: dict) -> bytes | None:
    """Turn a harvested `<img>` back into bytes, honouring the crop `object-fit: cover` applied.

    The harvester already reports the *painted* rect rather than the element box, so the aspect
    ratio is right by construction — non-uniform stretch is unrepresentable here, which is why
    the old OOXML aspect-ratio audit has nothing left to check.
    """
    m = _DATA_URI.match(node.get("src") or "")
    if not m:
        return None
    raw = base64.b64decode(m.group(3)) if m.group(2) else m.group(3).encode("utf-8")
    crop = node.get("crop")
    if not crop or Image is None:
        return raw
    try:
        im = Image.open(io.BytesIO(raw))
        w, h = im.size
        box = (int(w * crop["l"]), int(h * crop["t"]),
               int(w * (1 - crop["r"])), int(h * (1 - crop["b"])))
        if box[2] - box[0] < 1 or box[3] - box[1] < 1:
            return raw
        buf = io.BytesIO()
        im.crop(box).convert("RGBA").save(buf, "PNG")
        return buf.getvalue()
    except Exception:
        return raw


# ---------------------------------------------------------------------------------------------
# Emit
# ---------------------------------------------------------------------------------------------

_ALIGN = {"left": PP_ALIGN.LEFT, "start": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
          "right": PP_ALIGN.RIGHT, "end": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}

# What may legally become an external hyperlink relationship.
_EXTERNAL = re.compile(r"^(https?|mailto|ftp|file):", re.I)

_JUMP = "ppaction://hlinksldjump"
_SLIDE_RT = ("http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")


def _slide_jump(run, index: int) -> None:
    """Turn an in-deck anchor into a real "go to slide" click.

    The section-agenda's roadmap rows are the deck's one piece of genuine navigation, and they
    are anchors (`#sec-N`) that a `.pptx` cannot follow as URLs. PowerPoint's own equivalent is a
    slide-jump action, which needs an *internal* relationship to the target slide part — so the
    link keeps working instead of being dropped.

    Deferred: the target slide may not exist yet when the run is written, so the jump is recorded
    and wired up once every slide is in the deck (`_wire_jumps`).
    """
    _PENDING.append((run, index))


_PENDING: list = []


_NOTES_MASTER_RT = ("http://schemas.openxmlformats.org/officeDocument/2006/"
                    "relationships/notesMaster")


def _declare_notes_master(prs) -> bool:
    """Declare the notes master in `presentation.xml`, which python-pptx does not.

    Asking for a slide's `notes_slide` makes python-pptx create the notesMaster part and relate
    the presentation to it — but it never adds the matching `<p:notesMasterIdLst>` to
    `presentation.xml`. The part is present, the relationship is present, and the presentation
    never says the master exists.

    PowerPoint shrugs at that. **Keynote refuses the file outright** — and refuses the whole deck,
    not the notes: a 86-slide export died on it, and the failure was invisible from this side
    because the package is well-formed, python-pptx reads it back happily, and macOS's own Office
    preview renders it. It reproduced only as "the file format is invalid", and only in Keynote.

    Speaker notes are load-bearing here, so dropping them was never the fix.
    """
    el = prs._element
    if el.find(qn("p:notesMasterIdLst")) is not None:
        return False
    rid = next((r_id for r_id, rel in prs.part.rels.items()
                if rel.reltype == _NOTES_MASTER_RT), None)
    if rid is None:
        return False
    lst = el.makeelement(qn("p:notesMasterIdLst"), {})
    lst.append(el.makeelement(qn("p:notesMasterId"), {qn("r:id"): rid}))
    # Schema order: sldMasterIdLst, notesMasterIdLst, handoutMasterIdLst, sldIdLst, sldSz, …
    masters = el.find(qn("p:sldMasterIdLst"))
    if masters is not None:
        masters.addnext(lst)
    else:
        el.insert(0, lst)
    return True


def _wire_jumps(prs) -> int:
    slides = list(prs.slides)
    wired = 0
    for run, index in _PENDING:
        if not 0 <= index < len(slides):
            continue
        try:
            part = run.part
            rid = part.relate_to(slides[index].part, _SLIDE_RT)
            rPr = run.font._rPr
            link = rPr.makeelement(qn("a:hlinkClick"),
                                   {qn("r:id"): rid, "action": _JUMP})
            # `a:hlinkClick` is the last child of `a:rPr` in the schema sequence.
            rPr.append(link)
            wired += 1
        except Exception:
            continue
    _PENDING.clear()
    return wired


def add_box(slide, n: dict) -> None:
    w, h = n["w"], n["h"]
    if w < 0.5 or h < 0.5:
        return
    radius = max(n.get("radius") or [0])
    small = min(w, h)
    if radius >= small / 2 - 0.5 and abs(w - h) < max(2.0, small * 0.25):
        shape = MSO_SHAPE.OVAL
    elif radius > 0.5:
        shape = MSO_SHAPE.ROUNDED_RECTANGLE
    else:
        shape = MSO_SHAPE.RECTANGLE

    sp = slide.shapes.add_shape(shape, E(n["x"]), E(n["y"]), E(w), E(h))
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = max(0.0, min(0.5, radius / small))
        except Exception:
            pass
    # A stock autoshape arrives with a theme fill, a theme outline and a preset shadow. Every one
    # of those is a lie about the deck, so all three are cleared and only what was measured is
    # put back.
    sp.shadow.inherit = False
    fill = n.get("fill")
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = rgb(fill)
        _set_alpha(sp.fill.fore_color._xFill.find(qn("a:srgbClr")).getparent(),
                   alpha_of(fill, n.get("op", 1)))
    else:
        sp.fill.background()
    sp.line.fill.background()
    sp.text_frame.word_wrap = False

    bd = n.get("borders")
    if bd:
        sides = [v for v in bd.values() if v]
        uniform = (len(sides) == 4 and
                   all(abs(s["w"] - sides[0]["w"]) < 0.1 and s["c"] == sides[0]["c"] for s in sides))
        if uniform:
            sp.line.color.rgb = rgb(sides[0]["c"])
            sp.line.width = P(sides[0]["w"])
        else:
            # An accent band on one edge is the deck's dominant motif. Drawing a full outline
            # would invent three edges that are not there, so each visible side becomes its own
            # thin filled rect on the exact geometry that was measured.
            for side, v in bd.items():
                if not v:
                    continue
                if side == "left":
                    r = (n["x"], n["y"], v["w"], h)
                elif side == "right":
                    r = (n["x"] + w - v["w"], n["y"], v["w"], h)
                elif side == "top":
                    r = (n["x"], n["y"], w, v["w"])
                else:
                    r = (n["x"], n["y"] + h - v["w"], w, v["w"])
                b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(r[0]), E(r[1]), E(r[2]), E(r[3]))
                b.shadow.inherit = False
                b.fill.solid()
                b.fill.fore_color.rgb = rgb(v["c"])
                b.line.fill.background()
                b.text_frame.word_wrap = False


def add_text(slide, n: dict) -> None:
    paras = [p for p in n.get("paras") or [] if p or len(n["paras"]) > 1]
    if not any(any(r["s"].strip() for r in p) for p in paras):
        return
    sc = n.get("sc") or 1.0
    pad = n.get("pad") or [0, 0, 0, 0]

    # A block that renders on one line cannot re-break, so wrapping is turned off for it — the
    # tightly fitted type (titles, pills, stat numbers) is then immune to any residual difference
    # between Chrome's line breaker and PowerPoint's. Wrapped prose gets a hair of extra width
    # and grows into the free space the harvest measured below it, so an extra line lands in
    # whitespace instead of on the next card.
    wrap = bool(n.get("wrap"))
    w = n["w"] * 1.015 + 2 if wrap else n["w"] + 4
    lh = n.get("lh") or (n["fs"] * 1.25)
    h = n["h"] + (min(n.get("slack") or 0, lh) if wrap else 0)

    tb = slide.shapes.add_textbox(E(n["x"]), E(n["y"]), E(w), E(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left, tf.margin_top = E(pad[3]), E(pad[0])
    tf.margin_right, tf.margin_bottom = E(pad[1]), E(pad[2])
    if n.get("vertical"):
        tf._bodyPr.set("vert", "vert270")

    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = _ALIGN.get(n.get("align"), PP_ALIGN.LEFT)
        # Exact leading in points, not a multiple: the vertical rhythm then survives regardless
        # of how PowerPoint would otherwise compute "single" spacing for this face.
        p.line_spacing = P(lh * sc)
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        for r in para:
            if not r["s"]:
                continue
            m = r["m"]
            run = p.add_run()
            run.text = r["s"]
            f = run.font
            f.name = m.get("fam") or "Arial"
            f.size = P(max(1.0, (m.get("px") or n["fs"]) * sc))
            f.bold = bool(m.get("b"))
            f.italic = bool(m.get("i"))
            f.underline = bool(m.get("u"))
            col = m.get("col")
            if col:
                f.color.rgb = rgb(col)
            # Strike-through and letter-spacing have no python-pptx API but are load-bearing
            # here: the deck strikes text for "don't do this" and tracks its uppercase labels.
            if m.get("s"):
                f._rPr.set("strike", "sngStrike")
            if m.get("spc"):
                f._rPr.set("spc", str(int(round(m["spc"] * sc * 100 * PT_PX))))
            # Only a real external scheme may become a hyperlink relationship. A fragment or a
            # bare path produces a relationship with a non-URI target, and PowerPoint rejects
            # the whole package as "file format is invalid" — one dead link on one slide taking
            # the entire deck down with it.
            href = m.get("href") or ""
            if href and _EXTERNAL.match(href):
                try:
                    run.hyperlink.address = href
                except Exception:
                    pass
            elif m.get("jump") is not None:
                _slide_jump(run, int(m["jump"]))


def add_picture(slide, n: dict, blob: bytes) -> None:
    if not blob or n["w"] < 0.5 or n["h"] < 0.5:
        return
    try:
        slide.shapes.add_picture(io.BytesIO(blob), E(n["x"]), E(n["y"]), E(n["w"]), E(n["h"]))
    except Exception as e:  # a picture never silently disappears
        sys.stderr.write("[pptx] picture failed (%s)\n" % e)


def build(data: dict, workdir: Path, out: Path) -> tuple[int, list[str]]:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(SLIDE_W), Emu(SLIDE_H)
    # python-pptx's default template declares `type="screen4x3"`, and setting the dimensions
    # does not update it — leaving the package claiming one aspect while measuring another.
    prs._element.find(qn("p:sldSz")).set("type", "screen16x9")
    blank = prs.slide_layouts[6]

    # Every SVG is rasterized at the largest size any slide paints it, so a diagram reused small
    # and large is sharp in both places.
    sizes: dict[str, tuple[float, float]] = {}
    for s in data["slides"]:
        for n in s["nodes"]:
            if n["k"] == "svg":
                w, h = sizes.get(n["id"], (0, 0))
                sizes[n["id"]] = (max(w, n["w"]), max(h, n["h"]))
    rasters = rasterize_svgs(data.get("svgs") or {}, sizes, workdir)

    warnings: list[str] = []
    for s in data["slides"]:
        slide = prs.slides.add_slide(blank)
        if s.get("bg"):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = rgb(s["bg"])
        for n in s["nodes"]:
            k = n["k"]
            if k == "box":
                add_box(slide, n)
            elif k == "text":
                add_text(slide, n)
            elif k == "svg":
                add_picture(slide, n, rasters.get(n["id"]))
            elif k == "img":
                add_picture(slide, n, decode_img(n))
            elif k == "media":
                add_box(slide, {**n, "fill": {"r": 24, "g": 24, "b": 28, "a": 1},
                                "radius": [12, 12, 12, 12], "borders": None, "op": 1})
        if s.get("notes"):
            slide.notes_slide.notes_text_frame.text = s["notes"]
        for w in s.get("warnings") or []:
            warnings.append("slide %d: %s" % (s["i"] + 1, w))

    if _declare_notes_master(prs):
        sys.stderr.write("[pptx] declared the notes master in the presentation\n")

    jumps = _wire_jumps(prs)
    if jumps:
        sys.stderr.write("[pptx] %d roadmap links wired as slide jumps\n" % jumps)

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    _verify_package(out)
    return len(data["slides"]), warnings


def _verify_package(path: Path) -> None:
    """Refuse to hand over a file PowerPoint will not open.

    Everything here is well-formed XML that python-pptx accepts and that opens fine when read
    back with python-pptx — and still gets rejected at the package level. That gap is exactly
    where a broken export hides, so the check runs on the saved file rather than on the model.
    """
    import zipfile
    with zipfile.ZipFile(path) as z:
        pres = z.read("ppt/presentation.xml").decode("utf-8", "replace")
        has_part = any(n.startswith("ppt/notesMasters/") for n in z.namelist())
        if has_part and "notesMasterIdLst" not in pres:
            raise RuntimeError(
                "the package carries a notes master that presentation.xml never declares — "
                "Keynote refuses the whole file for this. (See _declare_notes_master.)")
        for name in z.namelist():
            if not name.endswith(".rels"):
                continue
            body = z.read(name).decode("utf-8", "replace")
            for m in re.finditer(r'Target="([^"]+)"\s+TargetMode="External"', body):
                if not _EXTERNAL.match(m.group(1)):
                    raise RuntimeError(
                        "%s links to %r as an external target, which is not a URI — PowerPoint "
                        "refuses the whole package for this. (Fix the hyperlink handling in "
                        "add_text.)" % (name, m.group(1)))


# ---------------------------------------------------------------------------------------------

def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description="Derive an editable .pptx from the rendered HTML deck.")
    ap.add_argument("--talk", help="Talk root (talks/<Talk>); deck = output/html/index.html")
    ap.add_argument("--deck", help="an HTML deck to export directly")
    ap.add_argument("-o", "--output", help="output .pptx (default: <talk>/output/final.pptx)")
    ap.add_argument("--style", default="", help="deck skin to pin (default: the deck's own)")
    ap.add_argument("--theme", default="", help="light | dark")
    ap.add_argument("--timeout", type=int, default=180)
    ap.add_argument("--dump-harvest", help="also write the raw display list here (debugging)")
    a = ap.parse_args(argv)

    if a.deck:
        deck = Path(a.deck)
        out = Path(a.output) if a.output else deck.parent / "final.pptx"
    elif a.talk:
        deck = Path(a.talk) / "output" / "html" / "index.html"
        out = Path(a.output) if a.output else Path(a.talk) / "output" / "final.pptx"
    else:
        ap.error("one of --talk or --deck is required")

    if not deck.is_file():
        sys.stderr.write("[pptx] no deck at %s — render the HTML first (build_html.py)\n" % deck)
        return 2

    workdir = out.parent / ".export"
    workdir.mkdir(parents=True, exist_ok=True)
    try:
        sys.stderr.write("[pptx] measuring the rendered deck…\n")
        data = harvest(deck, a.style, a.theme, a.timeout)
        if a.dump_harvest:
            Path(a.dump_harvest).write_text(json.dumps(data), encoding="utf-8")
        sys.stderr.write("[pptx] %d slides, %d distinct vector graphics\n"
                         % (len(data["slides"]), len(data.get("svgs") or {})))
        n, warnings = build(data, workdir, out)
    except (_chrome.ChromeMissing, RuntimeError) as e:
        sys.stderr.write("[pptx] %s\n" % e)
        return 2

    for w in warnings[:20]:
        sys.stderr.write("[pptx] %s\n" % w)
    sys.stderr.write("[pptx] %d slides → %s\n" % (n, out))
    print(out)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
