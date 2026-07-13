"""Build the free-form deck's fixed slides — cover + agenda — from metadata.

Free-form is **template-free**: unlike strict (which opens `base-template.pptx`),
free-form starts from a fresh `python-pptx` `Presentation()` — which ships a full
default Office theme + slideMaster, so the deck imports into Keynote — and builds its
two fixed slides (the contractually-fixed cover, and the agenda) *by code from
metadata*, not by substituting into a binary template. The renderer then designs
slides 2+ into the same deck.

This is the committed builder so the cover/agenda are reproducible and never
hand-authored per run (the exact §2 geometry lives here, once). Reads the Talk's
`final.md` frontmatter (`presentation`, `class`, `presenter`, `date`) for the cover and
its numbered `# N. <Section>` H1s for the agenda.

Usage:
    python3 freeform_deck.py --final talks/<Talk>/final.md -o <deck.pptx> \
        [--logo <png>] [--language es]

Requires python-pptx. The renderer opens the produced deck and adds the content slides.
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

# free-form §2 cover geometry (EMU) — the contractually-fixed cover.
CANVAS_W, CANVAS_H = 9144000, 5143500                       # 16:9
_TITLE   = dict(off=(496119, 536823),  ext=(8151763, 1948458), face="Helvetica Bold", pt=40.5, hex="1F1E1E", lnpct=104)
_CLASS   = dict(off=(496119, 2677269), ext=(6216923, 235297),  face="Helvetica Bold", pt=14.5, hex="1F1E1E")
_AUTHOR  = dict(off=(496119, 3219748), ext=(3295799, 560933),  face="Helvetica",      pt=11.5, hex="3B3535")
_LOGO    = dict(off=(7183562, 3248546), ext=(1469008, 1214065))

# Agenda geometry.
_AG_TITLE = dict(off=(496119, 536823), ext=(8151763, 700000), face="Helvetica Bold", pt=28, hex="1F1E1E")
_AG_ROW_X, _AG_ROW_Y0, _AG_ROW_W, _AG_ROW_H, _AG_ROW_STRIDE = 496119, 1500000, 8151763, 360000, 520000
_AG_ACCENT = "DA1B2E"   # active/first-row accent

_SECTION_RE = re.compile(r"^#\s+(\d+)\s*[.\-—]\s*(.+?)\s*$")
_FRONTMATTER_RE = re.compile(r"\A---\s*\n(.*?)\n---\s*\n", re.DOTALL)

_MONTHS = {
    "es": ["enero", "febrero", "marzo", "abril", "mayo", "junio", "julio",
           "agosto", "septiembre", "octubre", "noviembre", "diciembre"],
    "en": ["January", "February", "March", "April", "May", "June", "July",
           "August", "September", "October", "November", "December"],
}
_LABELS = {"es": ("Autor:", "Última Modificación:"), "en": ("Author:", "Last modified:")}


def _parse_frontmatter(text: str) -> dict:
    m = _FRONTMATTER_RE.match(text)
    fields: dict[str, str] = {}
    if not m:
        return fields
    for line in m.group(1).splitlines():
        if line.lstrip().startswith("#") or ":" not in line:
            continue
        k, _, v = line.partition(":")
        fields[k.strip()] = v.strip().strip('"').strip("'")
    return fields


def _sections(text: str) -> list[str]:
    out = []
    for line in text.splitlines():
        m = _SECTION_RE.match(line)
        if m:
            out.append(m.group(2).strip())
    return out


def _fmt_date(iso: str, lang: str) -> str:
    m = re.match(r"^\s*(\d{4})-(\d{2})-(\d{2})", iso or "")
    if not m:
        return iso or ""
    y, mo = int(m.group(1)), int(m.group(2))
    months = _MONTHS.get(lang, _MONTHS["es"])
    return f"{months[mo - 1]}, {y}" if 1 <= mo <= 12 else f"{y}"


def _textbox(slide, spec, runs, Emu, Pt, RGBColor, PP_ALIGN):
    tb = slide.shapes.add_textbox(Emu(spec["off"][0]), Emu(spec["off"][1]),
                                  Emu(spec["ext"][0]), Emu(spec["ext"][1]))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, hex_override) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = spec["face"]
        f.size = Pt(spec["pt"])
        f.color.rgb = RGBColor.from_string(hex_override or spec["hex"])
    return tb


def build(final_md: str, logo: Path | None, lang: str):
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    fm = _parse_frontmatter(final_md)
    sections = _sections(final_md)
    author_lbl, mod_lbl = _LABELS.get(lang, _LABELS["es"])

    prs = Presentation()
    prs.slide_width = Emu(CANVAS_W)
    prs.slide_height = Emu(CANVAS_H)
    blank = prs.slide_layouts[6]

    # ── Cover (§2) ──
    cov = prs.slides.add_slide(blank)
    _textbox(cov, _TITLE, [(fm.get("presentation", ""), None)], Emu, Pt, RGBColor, PP_ALIGN)
    _textbox(cov, _CLASS, [(fm.get("class", ""), None)], Emu, Pt, RGBColor, PP_ALIGN)
    _textbox(cov, _AUTHOR, [
        (f"{author_lbl} {fm.get('presenter', '')}".strip(), None),
        (f"{mod_lbl} {_fmt_date(fm.get('date', ''), lang)}".strip(), None),
    ], Emu, Pt, RGBColor, PP_ALIGN)
    if logo and logo.is_file():
        cov.shapes.add_picture(str(logo), Emu(_LOGO["off"][0]), Emu(_LOGO["off"][1]),
                               Emu(_LOGO["ext"][0]), Emu(_LOGO["ext"][1]))

    # ── Agenda ──
    ag = prs.slides.add_slide(blank)
    _textbox(ag, _AG_TITLE, [("Agenda", None)], Emu, Pt, RGBColor, PP_ALIGN)
    y = _AG_ROW_Y0
    for i, name in enumerate(sections, 1):
        spec = dict(off=(_AG_ROW_X, y), ext=(_AG_ROW_W, _AG_ROW_H), face="Helvetica", pt=16, hex="1F1E1E")
        _textbox(ag, spec, [(f"{i}.  {name}", _AG_ACCENT if i == 1 else None)], Emu, Pt, RGBColor, PP_ALIGN)
        y += _AG_ROW_STRIDE

    return prs, len(sections)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--final", type=Path, required=True, help="Path to talks/<Talk>/final.md")
    ap.add_argument("-o", "--output", type=Path, required=True, help="Deck to write (cover + agenda).")
    ap.add_argument("--logo", type=Path,
                    default=Path(__file__).resolve().parent.parent.parent
                    / "config" / "pptx-styles" / "free-form" / "cover-logo.png",
                    help="Cover logo PNG (default: the bundled free-form cover-logo.png).")
    ap.add_argument("--language", default="es", help="Cover label language (es|en; default es).")
    args = ap.parse_args()

    if not args.final.is_file():
        print(f"error: {args.final} not found", file=sys.stderr)
        return 2
    try:
        import pptx  # noqa: F401
    except ImportError:
        print("error: python-pptx is required (`pip install python-pptx`)", file=sys.stderr)
        return 3

    prs, n_sections = build(args.final.read_text(encoding="utf-8"), args.logo, args.language)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.output)
    print(f"[freeform] cover + agenda ({n_sections} sections) → {args.output}", file=sys.stderr)
    print(str(args.output))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
