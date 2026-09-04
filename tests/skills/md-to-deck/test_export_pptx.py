#!/usr/bin/env python3
"""Regression tests for the HTML -> PPTX export.

Run:  python3 tests/skills/md-to-deck/test_export_pptx.py

## Why this file exists

The `.pptx` is derived by *measuring* the rendered HTML deck: Chrome lays it out, `harvest.js`
reports a display list, `export_pptx.py` rebuilds it as native shapes. Nothing in that chain
fails loudly on its own — a harvester that silently stops registering, a font token that stops
reaching a text surface, an emitter that drops every picture on one slide: all of them produce a
`.pptx` that opens fine and is quietly wrong. These are the checks that make those failures loud.

The first two need no browser and run anywhere. The third does need Chrome and skips without it,
because the thing it guards (does the browser still hand back a complete harvest?) cannot be
tested any other way.
"""
from __future__ import annotations

import re
import sys
import tempfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[3]
SKILL = ROOT / "skills" / "md-to-deck"
FIXTURE = ROOT / "tests" / "skills" / "md-to-deck"
sys.path.insert(0, str(SKILL))

THEME = SKILL / "templates" / "html" / "theme.css"
HARVEST = SKILL / "templates" / "html" / "harvest.js"


def check_font_tokens() -> list[str]:
    """Every font-family in the deck must resolve through `--sans` / `--mono`.

    This guards the load-bearing claim behind the whole export: the geometry is harvested with
    the PPTX's own fonts already substituted in the browser, so what PowerPoint re-wraps matches
    what Chrome measured. A single hardcoded `font-family: Georgia` in theme.css would opt one
    surface out of that substitution, and the export would start lying about its own geometry
    with nothing anywhere to notice.
    """
    bad = []
    for i, line in enumerate(THEME.read_text(encoding="utf-8").splitlines(), 1):
        if "--sans:" in line or "--mono:" in line:      # the token declarations themselves
            continue
        for m in re.finditer(r"font-family\s*:\s*([^;}]+)", line):
            v = m.group(1).strip()
            if not (v.startswith("var(--sans") or v.startswith("var(--mono") or v == "inherit"):
                bad.append(f"theme.css:{i} font-family names a family directly: {v}")
    return bad


def check_harvest_source() -> list[str]:
    """The harvester is inlined into a <script> block, and such a block ends at the first closing
    script tag anywhere in its text — comments included. When that happened, the browser reported
    a syntax error on the leftover markup, the harvester never registered, and the export failed
    with no clue as to why. Cost three debugging rounds; costs one line to prevent."""
    src = HARVEST.read_text(encoding="utf-8")
    bad = []
    if "</scr" + "ipt" in src:
        bad.append("harvest.js contains a closing script tag — inlined, it truncates itself")
    import html_style
    if "</scr" + "ipt" in html_style.HARVEST_JS:
        bad.append("html_style.HARVEST_JS is not escaping the closing script tag")
    return bad


def check_package(pptx: Path) -> list[str]:
    """Package-level validity, checked on the saved file.

    This is the gap that actually bit: a deck whose XML is well-formed, that python-pptx reads
    back without complaint, and that PowerPoint still refuses with "the file format is invalid".
    The cause was the section-agenda roadmap rows — in-deck anchors (`#sec-N`) emitted as
    external hyperlink relationships, whose targets are not URIs. Seven slides' worth of dead
    links took the entire 86-slide deck down with them, and nothing upstream of PowerPoint said
    a word. So the invariants are asserted here, on the bytes that ship.
    """
    import posixpath
    import zipfile

    bad = []
    with zipfile.ZipFile(pptx) as z:
        names = set(z.namelist())
        if z.testzip() is not None:
            bad.append("the package is a corrupt zip")
        for n in z.namelist():
            if not n.endswith(".rels"):
                continue
            base = n[:n.rindex("/_rels/")] if "/_rels/" in n else ""
            body = z.read(n).decode("utf-8", "replace")
            for m in re.finditer(r'Target="([^"]+)"(\s+TargetMode="External")?', body):
                target, external = m.group(1), bool(m.group(2))
                if external:
                    if not re.match(r"^(https?|mailto|ftp|file):", target, re.I):
                        bad.append(f"{n}: external target {target!r} is not a URI")
                    continue
                cand = posixpath.normpath(posixpath.join(base, target)).lstrip("/")
                if cand not in names:
                    bad.append(f"{n}: relationship points at a missing part {cand}")
        pres = z.read("ppt/presentation.xml").decode("utf-8", "replace")
        m = re.search(r'<p:sldSz[^>]*cx="(\d+)"[^>]*cy="(\d+)"[^>]*/>', pres)
        if not m or (m.group(1), m.group(2)) != ("12192000", "6858000"):
            bad.append("presentation.xml does not declare a 16:9 slide")
        elif 'type="screen16x9"' not in m.group(0):
            bad.append("the slide size measures 16:9 but declares another aspect")
    return bad


def check_emitted_deck(pptx: Path, geom: dict) -> list[str]:
    """Structural parity between the display list and the emitted file."""
    from pptx import Presentation

    bad = []
    prs = Presentation(str(pptx))
    slides = list(prs.slides)

    if prs.slide_width != 12192000 or prs.slide_height != 6858000:
        bad.append(f"slide is {prs.slide_width}x{prs.slide_height} EMU, not 16:9 at 13.333x7.5in")
    if len(slides) != len(geom["slides"]):
        bad.append(f"{len(slides)} slides emitted for {len(geom['slides'])} harvested")

    fonts, notes, oob = set(), 0, 0
    for i, (slide, hs) in enumerate(zip(slides, geom["slides"]), 1):
        pics = sum(1 for sh in slide.shapes if sh.shape_type == 13)
        want = sum(1 for n in hs["nodes"] if n["k"] in ("svg", "img"))
        if pics != want:
            bad.append(f"slide {i}: {pics} pictures emitted for {want} harvested")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            notes += 1
        for sh in slide.shapes:
            if sh.left is None:
                continue
            if (sh.left < -2 * 9525 or sh.top < -2 * 9525
                    or sh.left + sh.width > prs.slide_width + 8 * 9525
                    or sh.top + sh.height > prs.slide_height + 8 * 9525):
                oob += 1
            if sh.has_text_frame:
                for p in sh.text_frame.paragraphs:
                    for r in p.runs:
                        fonts.add(r.font.name)

    want_notes = sum(1 for s in geom["slides"] if s.get("notes"))
    if notes != want_notes:
        bad.append(f"{notes} slides carry speaker notes, {want_notes} were harvested")
    # A `.pptx` cannot embed a font, so the export must name only faces every machine has.
    stray = {f for f in fonts if f not in ("Arial", "Courier New")}
    if stray:
        bad.append(f"fonts that will not resolve everywhere: {sorted(stray)}")
    if oob > 2:
        bad.append(f"{oob} shapes fall outside the slide")
    return bad


def main() -> int:
    failures = []

    print("font tokens own every text surface …")
    failures += check_font_tokens()
    print("the harvester survives being inlined …")
    failures += check_harvest_source()

    import _chrome
    deck = FIXTURE / "style-reference.html"
    if not _chrome.find():
        print("no Chrome found — skipping the end-to-end export check")
    elif not deck.is_file():
        print(f"no fixture deck at {deck} — skipping the end-to-end export check")
    else:
        print("exporting the reference deck …")
        import export_pptx
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "final.pptx"
            try:
                geom = export_pptx.harvest(deck)
                export_pptx.build(geom, Path(td), out)
                failures += check_package(out)
                failures += check_emitted_deck(out, geom)
            except Exception as e:
                failures.append(f"the export failed: {e}")

    if failures:
        print(f"\n{len(failures)} check(s) FAILED:")
        for f in failures:
            print(f"  - {f}")
        return 1
    print("\nall export checks pass.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
